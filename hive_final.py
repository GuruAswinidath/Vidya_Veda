import json
import pathlib
import re
from typing import List, Tuple, Optional
import pptx
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from sympy import content
# Constants
ICON_SIZE = Inches(0.8)
ICON_BG_SIZE = Inches(1)
ICON_COLORS = [
    RGBColor(41, 128, 185),  # Blue
    RGBColor(39, 174, 96),   # Green
    RGBColor(192, 57, 43),   # Red
    RGBColor(142, 68, 173),  # Purple
    RGBColor(243, 156, 18),  # Orange
    RGBColor(44, 62, 80)     # Dark Blue
]

PROGRESS_MARKER = '>>'
ICON_REGEX = re.compile(r'\[\[(.*?)\]\]\s*(.*)')
ICONS_DIR = 'icons\png128'  # Assuming icons are in an 'icons' directory
def add_comparison_table(slide, table_data, slide_width, slide_height):
    """
    Adds a comparison table to the slide based on the table data provided.
    :param slide: The slide where the table will be added.
    :param table_data: List of lists, where each sublist represents a row of the table.
    :param slide_width: Width of the slide.
    :param slide_height: Height of the slide.
    """
    n_rows = len(table_data)
    n_cols = len(table_data[0]) if n_rows > 0 else 0

    # Set table dimensions
    left = Inches(0.5)
    top = Inches(2)
    table_width = slide_width - Inches(1)  # Leave some margin
    table_height = Inches(3)  # Set a fixed height for the table

    # Add table shape
    table = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, table_height).table

    # Fill table data
    for row_idx, row in enumerate(table_data):
        for col_idx, cell_text in enumerate(row):
            table.cell(row_idx, col_idx).text = str(cell_text)

            # Formatting cells
            cell = table.cell(row_idx, col_idx)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if row_idx == 0:  # Header row formatting
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(41, 128, 185)  # Header background color
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Header font color
            else:
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Body font color

def generate_monochrome_slide(slide):
    """Sets a black background with white text for a monochrome look."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Update all text shapes to use white text
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text

def generate_powerpoint_presentation(parsed_data: dict, output_file_path: pathlib.Path):
    presentation = pptx.Presentation()
    slide_width = Inches(presentation.slide_width.inches)
    slide_height = Inches(presentation.slide_height.inches)

    # Title slide
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = parsed_data['title']
    title_slide.placeholders[1].text = "Created with VidyaVeda"
    generate_monochrome_slide(title_slide)

    for slide_data in parsed_data['slides']:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data['heading']
        generate_monochrome_slide(slide)  
        # Handle different slide contents based on data
        if any(isinstance(point, str) and PROGRESS_MARKER in point for point in slide_data['bullet_points']):
            add_progress_timeline(slide, slide_data['bullet_points'], slide_width, slide_height)
        elif any(isinstance(point, str) and ICON_REGEX.match(point) for point in slide_data['bullet_points']):
            add_icons_slide(slide, slide_data['bullet_points'], slide_width, slide_height)
        elif 'table_data' in slide_data and slide_data['table_data']:
            add_comparison_table(slide, slide_data['table_data'], slide_width, slide_height)
        else:
            add_default_content(slide, slide_data['bullet_points'])

        if 'key_message' in slide_data and slide_data['key_message']:
            add_key_message(slide, slide_data['key_message'], slide_width, slide_height)

        if 'img_keywords' in slide_data and slide_data['img_keywords']:
            add_image_placeholder(slide, slide_data['img_keywords'], slide_width, slide_height)

    presentation.save(output_file_path)
    print(f"Presentation saved as {output_file_path}")






def add_default_content(slide, bullet_points):
    # Define the custom width and position for the text box area
    left_margin = Inches(0.1)  # Left margin for the text box
    top_margin = Inches(1.5)   # Start lower down to give space for the title
    text_width = Inches(5.2)   # Width of the text box
    text_height = Inches(5)    # Height of the text box
    
    # Add a new text box with restricted width
    textbox = slide.shapes.add_textbox(left_margin, top_margin, text_width, text_height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True  # Enable word wrap to fit text within width

    for point in bullet_points:
        if isinstance(point, str):
            # Add each bullet point to the text frame
            p = text_frame.add_paragraph()
            p.bullet = True  
            p.text = point
            p.level = 0  # Top-level bullet
            p.font.size = Pt(16)  # Font size
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text for monochrome theme

            # Add a two-line space after each bullet point
            empty_p = text_frame.add_paragraph()
            empty_p.text = ""
            empty_p.level = 0
            empty_p.space_after = Pt(14)  # Two-line space (14pt * 2)
        elif isinstance(point, dict):
            # Add heading and sub-bullets if point is a dict with nested points
            p = text_frame.add_paragraph()
            p.bullet = True
            p.text = point['heading']
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            for subpoint in point['bullet_points']:
                sub_p = text_frame.add_paragraph()
                sub_p.text = subpoint
                sub_p.level = 2  # Indented bullet level
                sub_p.font.size = Pt(14)
                sub_p.font.color.rgb = RGBColor(200, 200, 200)
                sub_p.bullet = True

                # Add a two-line space after each sub-bullet point
                empty_sub_p = text_frame.add_paragraph()
                empty_sub_p.text = ""
                empty_sub_p.level = 2
                empty_sub_p.space_after = Pt(14)  # Two-line space (14pt * 2)

def add_progress_timeline(slide, bullet_points, slide_width, slide_height):
    steps = [point.strip(PROGRESS_MARKER).strip() for point in bullet_points if isinstance(point, str) and PROGRESS_MARKER in point]
    n_steps = len(steps)

    # Define the maximum number of blocks per row
    max_steps_per_row = 3
    n_rows = (n_steps + max_steps_per_row - 1) // max_steps_per_row  # Calculate the number of rows needed

    shape_width = Inches(slide_width.inches * 0.9 / min(n_steps, max_steps_per_row))
    shape_height = Inches(1.2)
    left_margin = Inches(0.5)
    top_margin = Inches(2.5)

    for row in range(n_rows):
        top = top_margin + row * (shape_height + Inches(0.5))  # Adjust the top position for each row
        left = left_margin
        
        # Determine how many steps should be placed in the current row
        steps_in_current_row = min(max_steps_per_row, n_steps - row * max_steps_per_row)
        
        for idx in range(steps_in_current_row):
            step_idx = row * max_steps_per_row + idx
            step = steps[step_idx]

            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, shape_width, shape_height)
            shape.text = step
            shape.text_frame.word_wrap = True
            shape.fill.solid()
            shape.fill.fore_color.rgb = ICON_COLORS[step_idx % len(ICON_COLORS)]
            left += shape_width + Inches(0.2)  # Adjust the position for the next step


import pathlib
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

def add_icons_slide(slide, bullet_points, slide_width, slide_height):
    icons_data = [ICON_REGEX.match(point).groups() for point in bullet_points if isinstance(point, str) and ICON_REGEX.match(point)]
    n_icons = len(icons_data)

    icon_size = ICON_SIZE
    spacing = Inches((slide_width.inches - n_icons * icon_size.inches) / (n_icons + 1))
    top = Inches(2)

    for idx, (icon_name, text) in enumerate(icons_data):
        left = spacing + idx * (icon_size + spacing)

        # Add icon background
        bg_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            left - Inches(0.1),
            top - Inches(0.1),
            ICON_BG_SIZE,
            ICON_BG_SIZE
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = ICON_COLORS[idx % len(ICON_COLORS)]

        # Add icon
        icon_path = pathlib.Path(ICONS_DIR) / f"{icon_name}.png"
        if icon_path.exists():
            slide.shapes.add_picture(str(icon_path), left, top, height=icon_size)
        else:
            placeholder = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                left,
                top,
                icon_size,
                icon_size
            )
            placeholder.text = icon_name

        # Add text below icon
        textbox = slide.shapes.add_textbox(
            left - Inches(0.3),
            top + icon_size + Inches(0.1),  # Align text directly below the icon
            icon_size + Inches(0.6),
            Inches(0.5)
        )
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)  # Change text color


def add_key_message(slide, message, slide_width, slide_height):
    left = Inches(0.5)
    top = slide_height - Inches(1.5)
    width = slide_width - Inches(1)
    height = Inches(0.8)

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 230, 230)
    shape.line.color.rgb = RGBColor(200, 200, 200)

    text_frame = shape.text_frame
    text_frame.text = f"Key Message: {message}"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

import requests
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image
import os
from io import BytesIO

def add_image_placeholder(slide, keywords, slide_width, slide_height):
    # Function to download the first image based on keywords from Bing
    def download_first_image(keywords):
        search_url = f"https://www.bing.com/images/search?q={keywords.replace(' ', '+')}&FORM=HDRSC2"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/85.0.4183.83 Safari/537.36",
            "Accept-Language": "en-US,en;q=0.9",
        }
        response = requests.get(search_url, headers=headers)
        soup = BeautifulSoup(response.text, 'html.parser')
        
        first_image = None
        try:
            img_tag = soup.find("a", {"class": "iusc"})  # First image link container
            if img_tag:
                img_json = img_tag["m"]
                img_url = img_json.split('"murl":"')[1].split('"')[0]  # Extract image URL
                first_image = img_url

                # Download the image with headers to avoid 403 error
                img_headers = {
                    "User-Agent": headers["User-Agent"],
                    "Referer": search_url  # Referer header to bypass 403 errors
                }
                img_data = requests.get(first_image, headers=img_headers).content
                img = Image.open(BytesIO(img_data))
                
                # Convert if the image is in WEBP format
                if img.format == 'WEBP':
                    img = img.convert("RGB")
                    img_path = f"{keywords.replace(' ', '_')}.jpg"
                    img.save(img_path, "JPEG")
                else:
                    img_path = f"{keywords.replace(' ', '_')}.{img.format.lower()}"
                    img.save(img_path)
                
                return img_path
        except Exception as e:
            print(f"Error downloading or converting image: {e}")
        
        return None

    # Define position for image/placeholder
    left = slide_width - Inches(4)
    top = Inches(2)
    width = Inches(3.5)
    height = Inches(4)

    # Attempt to download and place the image
    img_path = download_first_image(keywords)
    if img_path:
        # Add the image to the slide if download succeeds
        pic = slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.75), width=Inches(4.5), height=Inches(4.1))
        
        # Optionally delete the image after adding it to save space
        os.remove(img_path)
    else:
        # If no image found, use a placeholder
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray for contrast
        shape.line.color.rgb = RGBColor(200, 200, 200)
        
        text_frame = shape.text_frame
        text_frame.text = f"Image Placeholder\nKeywords: {keywords}"
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)  # Darker gray for readability




if __name__ == '__main__':
    with open('output.json', 'r') as json_file:
        data = json.load(json_file)

    output_path = pathlib.Path('output_presentation_2.pptx')
    generate_powerpoint_presentation(data, output_path)
    print(f"Presentation saved as {output_path}")
